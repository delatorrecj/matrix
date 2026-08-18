# -*- coding: utf-8 -*-
"""MATRIX Grand Finals deck (PPTX).

Projector fallback for `deck/index.html`. Same nine spoken slides, same six
appendix slides, same design system: the product's dark instrument surface,
the product's Geist type stack, one blue accent, five dimension colors.

Everything is authored against the same 1280x720 design canvas the HTML deck
uses, so a measurement here means the same thing there:

    u(px)     canvas pixels  ->  inches
    fs(px)    canvas pixels  ->  points

Speaker notes are written into PowerPoint's notes pane, so presenter view
carries the same cues as the HTML deck's S drawer.

Font: Geist is the product face. If it is not installed on the presenting
machine PowerPoint will substitute silently, so either install Geist or run
this with MATRIX_DECK_FONT="Segoe UI" for a predictable fallback.

    python build_deck.py
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "MATRIX_GrandFinals_PitchDeck.pptx")
ASSETS = os.path.join(HERE, "deck", "assets")

FONT = os.environ.get("MATRIX_DECK_FONT", "Geist")
FONT_MONO = os.environ.get("MATRIX_DECK_FONT_MONO", "Geist Mono")

# ── Palette, mirroring deck.css ──────────────────────────────────────────────
BG = RGBColor(0x0A, 0x0E, 0x1A)
BG_DEEP = RGBColor(0x06, 0x09, 0x0F)
SURFACE = RGBColor(0x11, 0x18, 0x27)
BORDER = RGBColor(0x1E, 0x29, 0x3B)
HAIRLINE = RGBColor(0x2C, 0x38, 0x4C)

TEXT = RGBColor(0xE9, 0xEE, 0xF8)
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)
TEXT_FAINT = RGBColor(0x7C, 0x8C, 0xA8)

ACCENT = RGBColor(0x4F, 0x82, 0xEE)
ACCENT_INK = RGBColor(0x93, 0xB4, 0xF7)
LIMIT = RGBColor(0xFB, 0xBF, 0x24)

DIMS = [
    ("Behavioral", RGBColor(0x60, 0xA5, 0xFA)),
    ("Social", RGBColor(0xF4, 0x72, 0xB6)),
    ("Economic", RGBColor(0xFA, 0xCC, 0x15)),
    ("Ecological", RGBColor(0x4A, 0xDE, 0x80)),
    ("Societal", RGBColor(0xC0, 0x84, 0xFC)),
]

# ── Canvas ───────────────────────────────────────────────────────────────────
SW_IN, SH_IN = 13.333, 7.5
CANVAS_W = 1280.0
PX = SW_IN / CANVAS_W          # one canvas pixel, in inches
MARGIN = 80                    # canvas px, matches .pad side padding

prs = Presentation()
prs.slide_width = Inches(SW_IN)
prs.slide_height = Inches(SH_IN)
BLANK = prs.slide_layouts[6]

CORE_TOTAL = 9
_core_seen = [0]


def u(px):
    """Canvas pixels to inches."""
    return Inches(px * PX)


def fs(px):
    """Canvas pixels to points (the canvas renders at 96 px per inch)."""
    return Pt(px * 0.75)


def _alpha(color_format, alpha):
    """python-pptx has no transparency API; set <a:alpha> on the fill directly."""
    srgb = color_format._xFill.find(qn("a:srgbClr"))
    node = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
    srgb.append(node)


def rect(s, x, y, w, h, color, alpha=None):
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, u(x), u(y), u(w), u(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if alpha is not None:
        _alpha(shape.fill.fore_color, alpha)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def scrim(s, x, y, w, h, a_start, a_end, angle, color=BG_DEEP):
    """A real gradient wash. Two stacked alpha rectangles leave a visible seam."""
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, u(x), u(y), u(w), u(h))
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.fill.gradient()
    shape.fill.gradient_angle = angle
    stops = shape.fill.gradient_stops
    for stop, pos, alpha in ((stops[0], 0.0, a_start), (stops[1], 1.0, a_end)):
        stop.position = pos
        stop.color.rgb = color
        _alpha(stop.color, alpha)
    return shape


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, CANVAS_W, 720, bg)
    return s


def text(s, x, y, w, h, paragraphs, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of dicts {runs: [(txt, px, weight, color, font?)], align, sa, sb, ls}"""
    box = s.shapes.add_textbox(u(x), u(y), u(w), u(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, pa in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        if pa.get("sb") is not None:
            p.space_before = fs(pa["sb"])
        if pa.get("sa") is not None:
            p.space_after = fs(pa["sa"])
        if pa.get("ls") is not None:
            p.line_spacing = pa["ls"]
        for run_spec in pa["runs"]:
            txt, size, bold, color = run_spec[:4]
            face = run_spec[4] if len(run_spec) > 4 else FONT
            spacing = run_spec[5] if len(run_spec) > 5 else None
            run = p.add_run()
            run.text = txt
            run.font.size = fs(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = face
            if spacing is not None:
                run.font._rPr.set("spc", str(int(spacing * 100)))
    return box


def hline(s, y, x=MARGIN, w=CANVAS_W - 2 * MARGIN, color=HAIRLINE, thick=1.2):
    return rect(s, x, y, w, thick, color)


def vline(s, x, y, h, color=HAIRLINE, thick=1.2):
    return rect(s, x, y, thick, h, color)


def footer(s, right=None, label="Team ATLAN", on_photo=False):
    ink = RGBColor(0xC3, 0xCD, 0xDF) if on_photo else TEXT_FAINT
    if right is None:
        _core_seen[0] += 1
        right = "%02d / %02d" % (_core_seen[0], CORE_TOTAL)
    text(
        s, MARGIN, 675, 500, 20,
        [{"runs": [("MATRIX", 12, False, ink, FONT_MONO, 0.06),
                   ("   |   ", 12, False, HAIRLINE, FONT_MONO, 0.06),
                   (label, 12, False, ink, FONT_MONO, 0.06)]}],
    )
    text(
        s, CANVAS_W - MARGIN - 300, 675, 300, 20,
        [{"runs": [(right, 12, False, ink, FONT_MONO, 0.06)], "align": PP_ALIGN.RIGHT}],
    )


def notes(s, body):
    s.notes_slide.notes_text_frame.text = body.strip()


def photo(s, filename, top=0.30, bottom=0.96):
    """Full-bleed image, cover-cropped, under one continuous vertical wash.
    A single full-width gradient keeps text at AA with no seam anywhere in
    frame; a partial-width side scrim leaves a visible vertical edge."""
    path = os.path.join(ASSETS, filename)
    if not os.path.isfile(path):
        print("  ! missing asset: %s (slide will render on flat background)" % filename)
        return False
    # Assets are 3:2. At full canvas width they stand 853 tall, so lift them
    # by half the overflow to cover 16:9 from the centre.
    s.shapes.add_picture(path, u(0), u(-66), width=u(CANVAS_W))
    scrim(s, 0, 0, CANVAS_W, 720, top, bottom, 270)
    return True


def wordmark(s, x=MARGIN, y=56, ink=TEXT):
    """Pentad mark plus wordmark: one kernel, five dimensions."""
    cx, cy, r = x + 13, y + 13, 9.5
    import math
    for i in range(5):
        a = math.radians(-90 + i * 72)
        px_, py_ = cx + r * math.cos(a), cy + r * math.sin(a)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, u(px_ - 2.2), u(py_ - 2.2), u(4.4), u(4.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ink
        dot.line.fill.background()
        dot.shadow.inherit = False
    core = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, u(cx - 3.5), u(cy - 3.5), u(7), u(7))
    core.fill.solid()
    core.fill.fore_color.rgb = ink
    core.line.fill.background()
    core.shadow.inherit = False
    text(s, x + 36, y + 3, 300, 26,
         [{"runs": [("MATRIX", 17, True, ink, FONT, 0.14)]}])


# ═════════════════════════════════════════════════════════════════════════════
# 1 · TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
photo(s, "hero-iloilo.png", top=0.30, bottom=0.94)
wordmark(s)
for name, xoff in (("logo_aaih.png", 0), ("logo_p2a_asean.png", 150)):
    p = os.path.join(ASSETS, name)
    if os.path.isfile(p):
        s.shapes.add_picture(p, u(CANVAS_W - MARGIN - 260 + xoff), u(52), height=u(30))

text(s, MARGIN, 330, 900, 210,
     [{"runs": [("See what a decision will do before a city has to live with it.", 54, True, TEXT)],
       "ls": 1.06}])
text(s, MARGIN, 566, 800, 60,
     [{"runs": [("Team ATLAN  ·  Polytechnic University of the Philippines", 14, False, TEXT_DIM, FONT_MONO)],
       "sa": 6},
      {"runs": [("ASEAN AI Hackathon 2026, Smart Cities track  ·  Pilot city: Iloilo", 14, False, TEXT_FAINT, FONT_MONO)]}])
footer(s, on_photo=True)
notes(s, """
COLD OPEN. Do not read the slide. Hold one beat, then open with the question.

SAY: "What if we could see the consequences of a decision before a city pays
the price for it?"

Then advance. Do not start with architecture. Do not start with "we built a
digital twin."
""")

# ═════════════════════════════════════════════════════════════════════════════
# 2 · THE QUESTION
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
text(s, MARGIN, 196, 980, 260,
     [{"runs": [("What if a city could see what a decision will do, before it has to live with the result?", 60, True, TEXT)],
       "ls": 1.06}])
text(s, MARGIN, 476, 760, 120,
     [{"runs": [("Today it cannot. Cities decide about roads, evacuation routes and transport without seeing what those decisions do to the people living there.", 24, False, TEXT_DIM)],
       "ls": 1.38}])
footer(s)
notes(s, """
BEAT 1, Cost of Not Knowing.

SAY: "Today, cities make decisions about roads, evacuation routes, transport
and infrastructure without being able to clearly see what those decisions will
do to the people living there."

Pause. Let the room sit in it. This beat is empathy, not product.
""")

# ═════════════════════════════════════════════════════════════════════════════
# 3 · COST OF NOT KNOWING
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
text(s, MARGIN, 64, 720, 120,
     [{"runs": [("Two problems sit underneath almost every bad urban decision.", 42, True, TEXT)], "ls": 1.12}])

ROOT_X, ROOT_W = MARGIN, 540
hline(s, 214, ROOT_X, ROOT_W)
text(s, ROOT_X, 234, ROOT_W, 90,
     [{"runs": [("Cities decide without seeing what happens next.", 27, False, TEXT)], "ls": 1.3}])
hline(s, 336, ROOT_X, ROOT_W)
text(s, ROOT_X, 356, ROOT_W, 110,
     [{"runs": [("The information they rely on is fragmented, inconsistent, or out of date.", 27, False, TEXT)], "ls": 1.3}])
hline(s, 478, ROOT_X, ROOT_W)

CAS_X, CAS_W = 720, 420
cascade = [
    "A road is built.",
    "Traffic shifts.",
    "An evacuation route becomes harder to reach.",
    "Flood risk moves to a barangay that was never in the plan.",
    "The city learns the cost after the concrete is poured.",
]
# A 420px column fits roughly 47 characters at 18px; budget rows accordingly.
cy = 214
for i, line in enumerate(cascade):
    last = i == len(cascade) - 1
    rect(s, CAS_X, cy + 11, 20 if last else 10, 1.2, ACCENT if last else HAIRLINE)
    text(s, CAS_X + 28, cy, CAS_W, 70,
         [{"runs": [(line, 18, last, TEXT if last else TEXT_DIM)], "ls": 1.4}])
    cy += 36 if len(line) <= 47 else 62
vline(s, CAS_X, 214, cy - 214 - 20)

hline(s, 540)
text(s, MARGIN, 564, 900, 90,
     [{"runs": [("Cities should not have to wait for reality to tell them they made the wrong decision.", 24, False, TEXT)],
       "ls": 1.32}])
footer(s)
notes(s, """
BEAT 1 continued.

DO NOT SAY "the Philippines has no data." The honest framing is fragmented,
inconsistent, or hard to operationalize.

Speak the cascade as one continuous image, not as five bullets. Then land:
"Cities shouldn't have to wait for reality to tell them they made the wrong
decision."

Target 45 to 55 seconds for beat 1 in total. Protected beat.
""")

# ═════════════════════════════════════════════════════════════════════════════
# 4 · WHAT IF
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
text(s, MARGIN, 176, 900, 180,
     [{"runs": [("“What happens if we put this here?”", 60, True, TEXT)], "ls": 1.06}])

hline(s, 372)
text(s, MARGIN, 400, 1000, 130,
     [{"runs": [("MATRIX", 27, True, TEXT),
                (" helps cities see the consequences of infrastructure decisions before those consequences become reality.", 27, False, TEXT)],
       "ls": 1.34}])
text(s, MARGIN, 552, 860, 80,
     [{"runs": [("Not another static study that ages the day it is filed. A way to explore possible futures, and to understand who pays if we are wrong.", 19, False, TEXT_FAINT)],
       "ls": 1.5}])
footer(s)
notes(s, """
BEAT 2, Promise land.

SAY: "So what if a city could ask, before it builds: what happens if we put
this here?"

Pause, then deliver the thesis sentence slowly. It is the single line you most
want the judges to repeat back.

Target 30 to 35 seconds. Protected beat.
""")

# ═════════════════════════════════════════════════════════════════════════════
# 5 · HOW IT WORKS
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
text(s, MARGIN, 64, 800, 80,
     [{"runs": [("Five steps, one simulated reality.", 42, True, TEXT)], "ls": 1.12}])

steps = [
    ("Simulate", "Plain language, or a pin on the map."),
    ("Visualize", "Watch the city move once that future exists."),
    ("Compare", "Five dimensions score the same run."),
    ("Identify", "Risks, and where our confidence is honestly low."),
    ("Act", "A brief a planner can answer."),
]
col_w = (CANVAS_W - 2 * MARGIN - 4 * 26) / 5.0
for i, (title, gloss) in enumerate(steps):
    x = MARGIN + i * (col_w + 26)
    rect(s, x, 208, col_w, 1.4, ACCENT if i == 0 else HAIRLINE)
    text(s, x, 226, col_w, 50,
         [{"runs": [(title, 28, True, TEXT)], "ls": 1.1}])
    text(s, x, 274, col_w, 130,
         [{"runs": [(gloss, 15, False, TEXT_DIM)], "ls": 1.42}])

hline(s, 470)
text(s, MARGIN, 494, 1120, 40,
     [{"runs": [("One kernel, one trajectory dataset, five scores that cannot contradict each other:", 16, False, TEXT_DIM)]}])
dx = MARGIN
for name, color in DIMS:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, u(dx), u(536), u(10), u(10))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    text(s, dx + 18, 530, 200, 30, [{"runs": [(name, 16, True, color)]}])
    dx += 22 + len(name) * 9.2
footer(s)
notes(s, """
BEAT 3, How it works.

Keep this beat plain. No SUMO, no Redis, no Azure. Those live in the appendix
and in Q&A.

The one line worth landing: "five dimensions, one shared simulated reality, so
the answers do not contradict each other."

COMPRESSIBLE BEAT. If time is short, say only "Simulate. Visualize. Act." and
move to the demo.
""")

# ═════════════════════════════════════════════════════════════════════════════
# 6 · DEMO
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
has_still = photo(s, "demo-still.png", top=0.34, bottom=0.92)
text(s, MARGIN, 272, 700, 120,
     [{"runs": [("Watch one decision.", 60, True, TEXT)], "ls": 1.06}])
beats_runs = []
for i, word in enumerate(["Scenario", "Simulation", "Consequence", "Decision"]):
    if i:
        beats_runs.append(("   ›   ", 24, False, ACCENT))
    beats_runs.append((word, 24, False, TEXT))
text(s, MARGIN, 414, 1000, 50, [{"runs": beats_runs}])
text(s, MARGIN, 486, 600, 40,
     [{"runs": [("matrix-atlan.vercel.app", 18, False, ACCENT_INK, FONT_MONO)]}])
rect(s, MARGIN, 516, 218, 1.2, ACCENT_INK, alpha=0.45)
footer(s, on_photo=has_still)
notes(s, """
BEAT 4, Demo. HARD PROTECTED BEAT.

20 to 30 seconds of video. Do not narrate the controls.

Scenario used: the Mandurriao school placement. See DEMO_SCRIPT.md for the live
fallback path.

If Inspect lands on screen: "That number isn't a guess. Equation, named
datasets, computed confidence. The AI narrates. It does not invent the figure."

To make this slide full-bleed, save a still as deck/assets/demo-still.png and
rerun build_deck.py.
""")

# ═════════════════════════════════════════════════════════════════════════════
# 7 · PROOF
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
text(s, MARGIN, 46, 900, 70,
     [{"runs": [("This is not just an idea.", 42, True, TEXT)], "ls": 1.12}])

ledger = [
    ("RUNS TODAY", "Unified simulation kernel, five impact modules, streaming API and Deck.gl frontend."),
    ("GROUNDED IN ILOILO", "Open-data foundation across 180 barangays and 5,680 priced parcels."),
    ("TRACEABLE", "Every scored number opens in Inspect: equation, named datasets, computed confidence."),
    ("GATED", "254 passing automated tests, plus two merge gates that block unprovenanced results."),
]
ly = 140
for label, value in ledger:
    hline(s, ly)
    text(s, MARGIN, ly + 20, 210, 40,
         [{"runs": [(label, 13, False, TEXT_FAINT, FONT_MONO, 0.08)]}])
    text(s, MARGIN + 238, ly + 14, CANVAS_W - 2 * MARGIN - 238, 70,
         [{"runs": [(value, 19, False, TEXT)], "ls": 1.35}])
    ly += 74
hline(s, ly)

LIM_Y = ly + 26
rect(s, MARGIN, LIM_Y, 2.4, 116, LIMIT)
text(s, MARGIN + 24, LIM_Y, 900, 30,
     [{"runs": [("And what we will not claim.", 18, True, LIMIT)]}])
limits = [
    "The behavioral validation gate is built. Its headline number is withheld until demand is calibrated.",
    "Flood back-testing is staged and not yet run against a real 2024 satellite extent.",
    "No city planner has formally signed off yet. That partnership is our next step.",
]
lyy = LIM_Y + 32
for line in limits:
    text(s, MARGIN + 24, lyy, 1050, 30,
         [{"runs": [(line, 16, False, TEXT_DIM)], "ls": 1.45}])
    lyy += 27

text(s, MARGIN, 616, 900, 40,
     [{"runs": [("False precision is the real risk. Honest confidence is the feature.", 20, False, TEXT)]}])
footer(s)
notes(s, """
BEAT 5, Proof.

Every figure here maps to CLAIMS.md. 254 = 190 kernel plus 64 API passing tests.

SAY: "Our empirical validation gates are built. We deliberately withhold the
behavioral headline until demand is calibrated, because publishing a confident
number off uncalibrated demand would break our own glass-box rule."

NEVER SAY: validated by CPDO, a published RMSE, or an invented accuracy
percentage.

Target 50 to 60 seconds.
""")

# ═════════════════════════════════════════════════════════════════════════════
# 8 · STAIRCASE
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
text(s, MARGIN, 64, 860, 120,
     [{"runs": [("We are not claiming the whole staircase. We are showing the first flight.", 42, True, TEXT)],
       "ls": 1.12}])

flights = [
    ("Simulation", "Explore what a decision would do.", 0.34),
    ("Data", "Connect to cleaner, continuously updated urban data.", 0.54),
    ("Digital twin", "A living model that learns from the city itself.", 0.76),
    ("Decision intelligence", "A regional way for cities to learn before they act.", 1.0),
]
BASE_Y, MAX_H = 620, 300
fw = (CANVAS_W - 2 * MARGIN - 3 * 20) / 4.0
for i, (title, body, frac) in enumerate(flights):
    x = MARGIN + i * (fw + 20)
    h = MAX_H * frac
    y = BASE_Y - h
    here = i == 0
    rect(s, x, y, fw, h, SURFACE if not here else RGBColor(0x14, 0x1E, 0x33))
    rect(s, x, y, fw, 1.6, ACCENT if here else HAIRLINE)
    rect(s, x, y, 1.6, h, ACCENT if here else HAIRLINE)
    text(s, x + 18, y + 18, fw - 36, 70,
         [{"runs": [(title, 25, True, TEXT)], "ls": 1.12}])
    # Titles longer than the column reserve two lines of headroom.
    text(s, x + 18, y + 18 + (74 if len(title) > 13 else 38), fw - 36, 110,
         [{"runs": [(body, 15, False, TEXT_DIM)], "ls": 1.4}])
    if here:
        # The shortest flight has no room inside it; the marker sits above.
        text(s, x, y - 28, fw, 24,
             [{"runs": [("WE ARE HERE", 12, False, ACCENT_INK, FONT_MONO, 0.08)]}])
footer(s)
notes(s, """
BEAT 6, Bigger vision.

The data gap is not an embarrassment. It is the reason the product has to climb.

SAY: "Today, MATRIX demonstrates what is possible through simulation. Tomorrow
it can connect to continuously updated urban data. Eventually, the city itself
becomes the feedback loop."

Steps 2, 3 and 4 are VISION, future tense only. Compressible beat, 30 to 40
seconds.
""")

# ═════════════════════════════════════════════════════════════════════════════
# 9 · THE CALL
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
photo(s, "iloilo-esplanade.png", top=0.28, bottom=0.96)
text(s, MARGIN, 356, 800, 100,
     [{"runs": [("Learn before they act.", 60, True, TEXT)], "ls": 1.06}])
text(s, MARGIN, 456, 700, 90,
     [{"runs": [("Every road, school and flood wall is a decision about someone’s life.", 26, False, RGBColor(0xD7, 0xDE, 0xEC))],
       "ls": 1.38}])
text(s, MARGIN, 570, 900, 40,
     [{"runs": [("Team ATLAN  ·  matrix-atlan.vercel.app  ·  Thank you", 14, False, TEXT_DIM, FONT_MONO)]}])
footer(s, on_photo=True)
notes(s, """
BEAT 7, the Call. HARD PROTECTED BEAT. Slow down. This is the last thing they
hear before scoring.

SAY: "We want Southeast Asian cities to stop learning from disasters,
congestion and failed infrastructure after the fact. We want them to learn
before they act. Because every infrastructure decision is ultimately a decision
about someone's life."

Then the soft ask: advance with us, from a working Iloilo pilot to the
decision-intelligence layer ASEAN cities deserve.
""")


# ═════════════════════════════════════════════════════════════════════════════
# APPENDIX. Q&A only, never the spoken spine.
# ═════════════════════════════════════════════════════════════════════════════
def appendix(tag, heading, heading_w=900):
    s = slide(BG_DEEP)
    text(s, MARGIN, 64, 400, 24,
         [{"runs": [("APPENDIX " + tag, 12, False, ACCENT_INK, FONT_MONO, 0.14)]}])
    text(s, MARGIN, 96, heading_w, 80,
         [{"runs": [(heading, 42, True, TEXT)], "ls": 1.12}])
    footer(s, right=tag, label="Appendix")
    return s


# A1 · Architecture
s = appendix("A1", "One kernel feeds five modules.")
nodes = [
    ("Orchestrator", "Azure OpenAI gpt-5.4 turns a natural-language ask or a map pin into a scenario. It never originates a number."),
    ("Simulation kernel", "SUMO via TraCI, a pre-warmed persona pool and a bias auditor, run as a delta against a nightly baseline."),
    ("One trajectory dataset", "A single per-agent record of the simulated city. Every dimension scores this same reality."),
    ("Synthesis and playback", "Citation-guarded narration, streamed over WebSocket into Deck.gl playback and the Inspect drawer."),
]
nw = (CANVAS_W - 2 * MARGIN - 3 * 34) / 4.0
for i, (title, body) in enumerate(nodes):
    x = MARGIN + i * (nw + 34)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, u(x), u(236), u(nw), u(174))
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    text(s, x + 16, 254, nw - 32, 40, [{"runs": [(title, 16, True, TEXT)], "ls": 1.15}])
    text(s, x + 16, 288, nw - 32, 120, [{"runs": [(body, 13, False, TEXT_DIM)], "ls": 1.4}])
    if i < 3:
        text(s, x + nw + 6, 310, 26, 30,
             [{"runs": [("›", 20, False, TEXT_FAINT)], "align": PP_ALIGN.CENTER}])

cw = (CANVAS_W - 2 * MARGIN - 4 * 16) / 5.0
sublabels = ["Travel time, mode shift", "Access and equity", "Land value, activity",
             "Emissions, exposure", "Services and resilience"]
for i, ((name, color), sub) in enumerate(zip(DIMS, sublabels)):
    x = MARGIN + i * (cw + 16)
    rect(s, x, 466, cw, 2.2, color)
    text(s, x, 480, cw, 30, [{"runs": [(name, 14, True, color)]}])
    text(s, x, 504, cw, 40, [{"runs": [(sub, 12.5, False, TEXT_DIM)], "ls": 1.35}])
notes(s, """
Pull this up only if a judge asks how it is built. The point of the diagram is
the middle box: one trajectory dataset is why five dimensions cannot contradict
each other.
""")

# A2 · Glass box
s = appendix("A2", "No number ships without its provenance.", heading_w=1000)
fields = [
    ("equation_id", "The exact equation from the methods ledger that produced the value. Locked, versioned, reviewable."),
    ("input_dataset_ids", "Every named dataset that fed the calculation, with its vintage and licence."),
    ("confidence", "Computed from data coverage and quality, never a guessed label. Low confidence suppresses the point estimate."),
]
fwid = (CANVAS_W - 2 * MARGIN - 2 * 22) / 3.0
for i, (code, body) in enumerate(fields):
    x = MARGIN + i * (fwid + 22)
    hline(s, 232, x, fwid)
    text(s, x, 252, fwid, 30, [{"runs": [(code, 15, False, ACCENT_INK, FONT_MONO)]}])
    text(s, x, 282, fwid, 140, [{"runs": [(body, 15, False, TEXT_DIM)], "ls": 1.45}])
text(s, MARGIN, 432, 980, 90,
     [{"runs": [("A merge gate enforces this in the build. If a result cannot resolve its equation, its datasets and a computed confidence in the Inspect drawer, it does not ship.", 19, False, TEXT)],
       "ls": 1.55}])
notes(s, """
This is the answer to "how do we know the AI is not making it up." The LLM
narrates and cites. It does not originate a figure.
""")

# A3 · Validation status
s = appendix("A3", "Validation, stated honestly.")
rows = [
    ("VAL-01", "Behavioral travel times against the Calderon Iloilo corridor study. Threshold is a\nnormalized RMSE of 0.30, per FHWA guidance.", "WITHHELD", LIMIT),
    ("VAL-02", "Flood extent against a real 2024 Sentinel-1 observation, scored by intersection over union.", "NOT RUN", LIMIT),
    ("VAL-03", "Mode share held within tolerance of the anchored city baseline.", "ENFORCED", DIMS[3][1]),
]
COL_A, COL_B, COL_C = MARGIN, MARGIN + 130, CANVAS_W - MARGIN - 130
rect(s, MARGIN, 210, CANVAS_W - 2 * MARGIN, 1.4, BORDER)
for h, x in (("GATE", COL_A), ("WHAT IT TESTS", COL_B), ("STATUS", COL_C)):
    text(s, x, 186, 300, 24, [{"runs": [(h, 12, False, TEXT_FAINT, FONT_MONO, 0.08)]}])
ry = 210
for gate, what, status, color in rows:
    text(s, COL_A, ry + 16, 120, 30, [{"runs": [(gate, 15, False, TEXT)]}])
    text(s, COL_B, ry + 16, COL_C - COL_B - 30, 70, [{"runs": [(what, 15, False, TEXT_DIM)], "ls": 1.4}])
    text(s, COL_C, ry + 16, 140, 30, [{"runs": [(status, 13, False, color, FONT_MONO, 0.04)]}])
    ry += 82 if "\n" in what else 62
    rect(s, MARGIN, ry, CANVAS_W - 2 * MARGIN, 1.2, HAIRLINE)
text(s, MARGIN, ry + 30, 1000, 100,
     [{"runs": [("The machinery for all three is implemented and tested. VAL-01 is withheld rather than published because demand volume is not yet calibrated, and a confident number off uncalibrated demand would break our own glass-box rule.", 19, False, TEXT)],
       "ls": 1.55}])
notes(s, """
"Withheld" is the correct word, not "failed" and not "not validated." The gate
exists and runs. We are choosing not to publish an uncalibrated headline.
""")

# A4 · Latency
s = appendix("A4", "Architected for a 90 second answer.")
lat = [
    ("TARGET", "90 seconds end to end, from question to a narrated five-dimension result."),
    ("HOW", "Pre-warmed persona pool, delta simulation against a nightly baseline, five modules in parallel, progressive streaming to the UI."),
    ("REPEAT RUNS", "Served from the trajectory cache, effectively instant."),
    ("COLD RUNS", "Can exceed the target. We quote it as an engineered target, not a measured guarantee."),
]
ly = 216
for label, body in lat:
    text(s, MARGIN, ly, 190, 30, [{"runs": [(label, 12, False, TEXT_FAINT, FONT_MONO, 0.07)]}])
    text(s, MARGIN + 200, ly - 4, CANVAS_W - 2 * MARGIN - 200, 70,
         [{"runs": [(body, 16, False, TEXT)], "ls": 1.45}])
    ly += 54
    rect(s, MARGIN, ly - 12, CANVAS_W - 2 * MARGIN, 1.2, HAIRLINE)
notes(s, """
Do not claim "always under 90 seconds." If you have a live measured figure from
the day, quote that figure instead.
""")

# A5 · Path
s = appendix("A5", "Iloilo is the beachhead, not the ceiling.")
legs = [
    ("Iloilo", "Working pilot on an open-data foundation, with a planner feedback loop built into the product.", True),
    ("Philippine cities", "Geographic scaling is an API-level change: swap the OSM bounding box and the local data layers.", False),
    ("ASEAN", "Behavioral scaling is prompt-level: reweight the persona archetypes to the city’s own mode share.", False),
]
lw = (CANVAS_W - 2 * MARGIN - 2 * 28) / 3.0
for i, (title, body, here) in enumerate(legs):
    x = MARGIN + i * (lw + 28)
    rect(s, x, 224, 1.6, 150, ACCENT if here else HAIRLINE)
    text(s, x + 22, 224, lw - 30, 44, [{"runs": [(title, 24, True, TEXT)], "ls": 1.12}])
    text(s, x + 22, 268, lw - 30, 120, [{"runs": [(body, 15, False, TEXT_DIM)], "ls": 1.45}])
text(s, MARGIN, 424, 900, 40,
     [{"runs": [("The engine is deliberately city-agnostic. Nothing in the kernel hard-codes Iloilo.", 19, False, TEXT)]}])
notes(s, """
This is a path, not a deployment claim. Never say MATRIX is already running
across ASEAN.
""")

# A6 · Stack
s = appendix("A6", "What it is made of.")
stack = [
    ("SIMULATION", "Eclipse SUMO via the TraCI Python API"),
    ("LANGUAGE MODEL", "Azure OpenAI gpt-5.4, orchestration and synthesis"),
    ("BACKEND", "FastAPI with a progressive WebSocket stream"),
    ("RETRIEVAL", "ChromaDB GraphRAG corpus, ingested at startup"),
    ("DATA", "Postgres with PostGIS, Redis for pools and caches"),
    ("FRONTEND", "Next.js 14 with Deck.gl animated playback"),
    ("DEPLOY", "Vercel for web, Hugging Face Spaces for the API"),
    ("FORECASTING", "XGBoost baseline forecaster"),
]
colw = (CANVAS_W - 2 * MARGIN - 56) / 2.0
for i, (label, body) in enumerate(stack):
    col, row = i % 2, i // 2
    x = MARGIN + col * (colw + 56)
    y = 216 + row * 62
    text(s, x, y, 160, 26, [{"runs": [(label, 12, False, TEXT_FAINT, FONT_MONO, 0.07)]}])
    text(s, x + 172, y - 4, colw - 172, 50, [{"runs": [(body, 15, False, TEXT)], "ls": 1.4}])
    rect(s, x, y + 40, colw, 1.2, HAIRLINE)
notes(s, """
Only open this if asked directly about the stack. It is a credibility card, not
part of the spoken arc.
""")


prs.save(OUT)
print("Wrote %s" % OUT)
_n = len(prs.slides._sldIdLst)
print("  %d slides (%d spoken + %d appendix), speaker notes in the notes pane."
      % (_n, CORE_TOTAL, _n - CORE_TOTAL))
if FONT == "Geist":
    print("  Font: Geist. If it is not installed here, PowerPoint substitutes silently.")
    print("  For a guaranteed fallback: MATRIX_DECK_FONT=\"Segoe UI\" python build_deck.py")
if not os.path.isfile(os.path.join(ASSETS, "demo-still.png")):
    print("  Demo slide is a cue card. Save deck/assets/demo-still.png and rerun for full bleed.")
sys.stdout.flush()
