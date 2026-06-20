# -*- coding: utf-8 -*-
"""Build Smart Cities_PUP_ATLAN_PitchDeck.pptx on the AAIH brand chrome.
Engine: python-pptx. Slide canvas 20in x 12.5in (matches AAIH template 16:10)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
A = lambda p: os.path.join(HERE, "_assets", p)

# ---------- palette ----------
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0xEA, 0xF1, 0xFF)  # near-white body
MUTE    = RGBColor(0x9D, 0xB2, 0xDA)  # muted caption
FAINT   = RGBColor(0x7E, 0x93, 0xBE)  # fainter
COBALT  = RGBColor(0x3B, 0x82, 0xF6)
CYAN    = RGBColor(0x38, 0xBD, 0xF8)
PANEL   = RGBColor(0x0B, 0x1E, 0x4D)  # card fill on dark bg
PANEL2  = RGBColor(0x0E, 0x24, 0x56)  # lighter card
BORDER  = RGBColor(0x2C, 0x50, 0xA6)  # card border
DARKBAR = RGBColor(0x06, 0x16, 0x3A)
# five dimensions
def C(h): return RGBColor((h >> 16) & 255, (h >> 8) & 255, h & 255)
BEH = C(0x3B82F6); SOC = C(0xF4476A); ECO = C(0xF59E0B)
ENV = C(0x22C55E); SOCI= C(0x8B5CF6)
GOOD= C(0x22C55E); WARN= C(0xF5A524); BAD = C(0xEF4444)
DIMS = [("Behavioral",BEH),("Social",SOC),("Economic",ECO),("Ecological",ENV),("Societal",SOCI)]

F_HEAD = "Segoe UI Semibold"
F_BODY = "Segoe UI"
F_NUM  = "Consolas"

SW, SH = 20.0, 12.5
prs = Presentation()
prs.slide_width  = Emu(18288000)
prs.slide_height = Emu(11430000)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, img="bg_dark.png"):
    s.shapes.add_picture(A(img), 0, 0, Inches(SW), Inches(SH))

def _set_radius(shape, r=0.09):
    try: shape.adjustments[0] = r
    except Exception: pass

def panel(s, x, y, w, h, fill=PANEL, line=BORDER, lw=1.0, radius=0.06, accent=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_radius(sp, radius)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if accent:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y+0.14), Inches(0.09), Inches(h-0.28))
        _set_radius(bar, 0.5)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit=False
    return sp

def rect(s, x, y, w, h, fill, line=None, lw=1.0, radius=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None: _set_radius(sp, radius)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def tb(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {runs:[(text,size,bold,color,font,italic)], align, sb, sa, sp}"""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, 0)
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        if pa.get("sb") is not None: p.space_before = Pt(pa["sb"])
        if pa.get("sa") is not None: p.space_after = Pt(pa["sa"])
        if pa.get("ls") is not None: p.line_spacing = pa["ls"]
        for r in pa["runs"]:
            text, size = r[0], r[1]
            bold  = r[2] if len(r) > 2 else False
            color = r[3] if len(r) > 3 else INK
            font  = r[4] if len(r) > 4 else F_BODY
            ital  = r[5] if len(r) > 5 else False
            run = p.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold; run.font.italic = ital
            run.font.name = font; run.font.color.rgb = color
    return box

def chip(s, x, y, w, h, text, fill, tcolor=WHITE, size=12, bold=True, line=None):
    sp = rect(s, x, y, w, h, fill, line=line, lw=1.0, radius=0.5)
    tf = sp.text_frame; tf.word_wrap = False
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,Inches(0.04))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size=Pt(size); r.font.bold=bold; r.font.name=F_HEAD; r.font.color.rgb=tcolor
    return sp

def arrow(s, x, y, w, h, color=COBALT):
    sp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit=False
    return sp

def band_and_logos(s):
    rect(s, 0, 0, SW, 1.04, WHITE)
    s.shapes.add_picture(A("logo_aaih.png"), Inches(0.7), Inches(0.17), height=Inches(0.72))
    s.shapes.add_picture(A("logo_p2a_asean.png"), Inches(17.92), Inches(0.20), height=Inches(0.66))

def header(s, kicker, title, page):
    band_and_logos(s)
    tb(s, 0.7, 1.30, 18.6, 0.4, [{"runs":[(kicker.upper(), 14.5, True, CYAN, F_HEAD)]}])
    tb(s, 0.7, 1.66, 18.6, 1.0, [{"runs":[(title, 33, True, WHITE, F_HEAD)]}])
    footer(s, page)

def footer(s, page):
    tb(s, 0.7, 11.95, 12.0, 0.4,
       [{"runs":[("MATRIX  ·  Pre-Construction Infrastructure Impact Simulator  ·  Team ATLAN", 11, False, FAINT, F_BODY)]}])
    tb(s, 17.3, 11.95, 2.0, 0.4,
       [{"runs":[("%02d / 15" % page, 11.5, True, MUTE, F_NUM)]}], )

def stat(s, x, y, w, num, label, color=CYAN, numsize=46):
    tb(s, x, y, w, 0.9, [{"runs":[(num, numsize, True, color, F_NUM)]}], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, x, y+0.86, w, 0.7, [{"runs":[(label, 12.5, False, MUTE, F_BODY)]}])

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = slide(); bg(s, "bg_title.png")
rect(s, 0, 0, SW, 1.04, WHITE)
s.shapes.add_picture(A("logo_aaih.png"), Inches(0.7), Inches(0.17), height=Inches(0.72))
s.shapes.add_picture(A("logo_p2a_asean.png"), Inches(17.92), Inches(0.20), height=Inches(0.66))
chip(s, 7.05, 2.55, 5.9, 0.62, "SMART CITIES  ·  ASEAN AI HACKATHON 2026", C(0x0A1E55), CYAN, size=15, line=COBALT)
tb(s, 1.0, 3.35, 18.0, 2.4, [{"runs":[("MATRIX", 150, True, WHITE, F_HEAD)], "align":PP_ALIGN.CENTER}])
tb(s, 1.0, 5.95, 18.0, 0.8,
   [{"runs":[("Multi-Agent Twin for Routing & Infrastructure eXchange", 30, False, INK, F_BODY)],"align":PP_ALIGN.CENTER}])
tb(s, 1.0, 7.05, 18.0, 0.8,
   [{"runs":[("See a project’s full impact — ", 25, False, MUTE, F_BODY),
             ("before you build it.", 25, True, CYAN, F_HEAD)],"align":PP_ALIGN.CENTER}])
# meta strip
panel(s, 4.6, 8.5, 10.8, 1.5, fill=C(0x081A45), line=BORDER, radius=0.10)
tb(s, 4.8, 8.72, 10.4, 1.2, [
   {"runs":[("Team ATLAN", 18, True, WHITE, F_HEAD),
            ("   ·   Polytechnic University of the Philippines", 18, False, INK, F_BODY)],"align":PP_ALIGN.CENTER,"sa":4},
   {"runs":[("Smart Cities Track    ·    Pilot City: Iloilo City, Philippines", 15, False, MUTE, F_BODY)],"align":PP_ALIGN.CENTER},
  ], anchor=MSO_ANCHOR.MIDDLE)
s.shapes.add_picture(A("logo_wordmark.png"), Inches(SW/2-1.6), Inches(11.75), height=Inches(0.42))

# =====================================================================
# SLIDE 2 — PROBLEM
# =====================================================================
s = slide(); bg(s); header(s, "The ASEAN Challenge", "Cities are building faster than they can foresee.", 2)
tb(s, 0.7, 2.78, 18.6, 0.9, [{"runs":[
    ("ASEAN is in its largest urban build-out in history — and infrastructure is still approved on "
     "static feasibility studies that age the day they are filed. The damage shows up on opening day.",
     17.5, False, INK, F_BODY)]}])
cards = [
  ("Studies age on filing", "A 2022 traffic count cannot predict how a 2026 development reshapes a corridor, "
   "a jeepney route, or a flood path.", ECO),
  ("Impacts judged in silos", "Environment, transport, economy, and equity are each reviewed in a different "
   "office — they meet only after the project is approved.", SOC),
  ("Tools need specialists", "Vissim, Aimsun, CityEngine, Replica — powerful, but none takes a planner’s "
   "plain-language question and answers across domains.", BEH),
]
cw, gap, x0, cy = 5.89, 0.46, 0.7, 4.0
for i,(t,b,c) in enumerate(cards):
    x = x0 + i*(cw+gap)
    panel(s, x, cy, cw, 3.05, accent=c, radius=0.05)
    tb(s, x+0.45, cy+0.32, cw-0.75, 0.7, [{"runs":[(t, 19, True, WHITE, F_HEAD)]}])
    tb(s, x+0.45, cy+1.18, cw-0.75, 1.7, [{"runs":[(b, 15, False, INK, F_BODY)],"ls":1.08}])
# bottom band: who pays
panel(s, 0.7, 7.45, 18.6, 1.55, fill=C(0x2A0E1E), line=C(0x7A2742), radius=0.06, accent=SOC)
tb(s, 1.15, 7.66, 17.9, 1.2, [
   {"runs":[("Who pays for the blind spot:  ", 16.5, True, SOC, F_HEAD),
            ("the commuter who loses 25–35% of the fare to first-mile travel no one modeled  ·  the informal "
             "vendor paved over because no one ran the displacement scenario  ·  the barangay that floods because "
             "runoff was calculated in isolation.", 16, False, INK, F_BODY)],"ls":1.06}])
tb(s, 0.7, 9.25, 18.6, 0.5, [{"runs":[
   ("Illustrative of a documented pattern (Montalbo, UP SURP; JICA 2014/2019; ICLEI Iloilo Roadmap) — "
    "shown as a pattern, not as cited point statistics.", 11.5, False, FAINT, F_BODY, True)]}])

# =====================================================================
# SLIDE 3 — THE AI SOLUTION (elevator pitch)
# =====================================================================
s = slide(); bg(s); header(s, "The AI Solution", "Drop a project on your city. See five futures in 90 seconds.", 3)
tb(s, 0.7, 2.78, 18.6, 1.0, [{"runs":[
   ("A planner asks in plain language — ", 17.5, False, INK, F_BODY),
   ("“what happens if we build an 8-storey mixed-use here?”", 17.5, True, CYAN, F_HEAD),
   ("  — or drops a pin. One unified kernel simulates thousands of agents; five impact modules score the "
    "same reality; the map animates and a cited narrative explains it.", 17.5, False, INK, F_BODY)]}])
# flow row
nodes = [("Plain-language\nquery / map drop", COBALT),
         ("Gemini orchestrator\nparses the scenario", COBALT),
         ("Unified kernel\nSUMO + AI personas", CYAN),
         ("Five impact modules\nscore one reality", CYAN),
         ("Animated map +\ncited report", COBALT)]
nw, ng, nx, ny = 3.22, 0.62, 0.7, 4.2
for i,(t,c) in enumerate(nodes):
    x = nx + i*(nw+ng)
    panel(s, x, ny, nw, 1.55, fill=PANEL2, line=c, lw=1.25, radius=0.07)
    tb(s, x+0.2, ny+0.2, nw-0.4, 1.2, [{"runs":[(t, 14.5, True, WHITE, F_HEAD)],"align":PP_ALIGN.CENTER,"ls":1.05}],
       anchor=MSO_ANCHOR.MIDDLE)
    if i < len(nodes)-1:
        arrow(s, x+nw+0.07, ny+0.55, 0.48, 0.45, COBALT)
# five dimension chips
tb(s, 0.7, 6.3, 18.6, 0.5, [{"runs":[("ONE SIMULATED REALITY, SCORED FIVE WAYS — EACH WITH AN EXPLICIT CONFIDENCE LEVEL", 13.5, True, MUTE, F_HEAD)]}])
dcw, dg, dx, dy = 3.48, 0.30, 0.7, 6.95
for i,(name,c) in enumerate(DIMS):
    x = dx + i*(dcw+dg)
    panel(s, x, dy, dcw, 1.15, fill=PANEL, line=c, lw=1.5, radius=0.10, accent=c)
    tb(s, x+0.35, dy+0.2, dcw-0.5, 0.8, [
       {"runs":[(name, 17, True, WHITE, F_HEAD)],"sa":2},
       {"runs":[("Impact dimension", 11.5, False, MUTE, F_BODY)]}], anchor=MSO_ANCHOR.MIDDLE)
panel(s, 0.7, 8.55, 18.6, 0.95, fill=C(0x081A45), line=BORDER, radius=0.08)
tb(s, 1.0, 8.7, 18.0, 0.7, [{"runs":[
   ("From  ", 17, False, MUTE, F_BODY), ("“is this allowed?”", 17, True, WHITE, F_HEAD),
   ("  to  ", 17, False, MUTE, F_BODY), ("“what will this actually do?”", 17, True, CYAN, F_HEAD),
   ("   No specialist. No modeling background. Ranges, not false precision.", 16, False, INK, F_BODY)]},],
   anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# SLIDE 4 — COMPETITIVE ADVANTAGE
# =====================================================================
s = slide(); bg(s); header(s, "Competitive Advantage", "The combination no existing tool offers.", 4)
rows = [
  ("Tool", "Plain-language input", "5 dimensions, one run", "Per-dimension confidence", "No specialist needed"),
  ("PTV Vissim / Aimsun", "—", "—", "—", "—"),
  ("ESRI CityEngine", "—", "—", "—", "—"),
  ("Replica / UrbanFootprint", "—", "partial", "—", "partial"),
  ("AnyLogic", "—", "partial", "—", "—"),
  ("MATRIX", "Yes", "Yes", "Yes", "Yes"),
]
tx, ty, tw = 0.7, 2.85, 12.4
rh = 0.78; colw = [3.4, 2.25, 2.25, 2.4, 2.1]
# header + rows
yy = ty
for ri, row in enumerate(rows):
    xx = tx
    is_head = ri == 0; is_matrix = row[0] == "MATRIX"
    fillc = COBALT if is_head else (C(0x10336F) if is_matrix else PANEL)
    linec = CYAN if is_matrix else BORDER
    for ci, cell in enumerate(row):
        w = colw[ci]
        rect(s, xx, yy, w, rh, fillc, line=linec, lw=1.25 if is_matrix else 0.75)
        if ci == 0:
            col = WHITE if (is_head or is_matrix) else INK
            tb(s, xx+0.2, yy, w-0.3, rh, [{"runs":[(cell, 14.5 if is_matrix else 13.5, True if (is_head or is_matrix) else False, col, F_HEAD if (is_head or is_matrix) else F_BODY)]}], anchor=MSO_ANCHOR.MIDDLE)
        else:
            if is_head:
                tb(s, xx+0.12, yy, w-0.2, rh, [{"runs":[(cell, 12, True, WHITE, F_HEAD)],"align":PP_ALIGN.CENTER,"ls":0.95}], anchor=MSO_ANCHOR.MIDDLE)
            else:
                mark = cell
                mc = GOOD if cell=="Yes" else (WARN if cell=="partial" else FAINT)
                disp = "✓" if cell=="Yes" else ("~" if cell=="partial" else "—")
                tb(s, xx+0.12, yy, w-0.2, rh, [{"runs":[(disp, 18 if cell=="Yes" else 15, True, mc, F_HEAD)],"align":PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        xx += w
    yy += rh
tb(s, tx, yy+0.12, tw, 0.5, [{"runs":[("Based on our competitor feature survey (GTM §2.1) — a combination claim, not an unfalsifiable absolute.", 11.5, False, FAINT, F_BODY, True)]}])
# right column: 3 UVPs
ux, uw = 13.5, 5.8
uvps = [("One kernel → five dimensions","All five score the same simulated reality, so they can never contradict each other.",CYAN),
        ("Glass box, by contract","Every number resolves to its equation, its datasets, and a computed confidence. The LLM narrates — it never invents a number.",COBALT),
        ("Plain language → 90 seconds","A pin or a sentence in; a calibrated, confidence-anchored five-dimension answer out.",BEH)]
uy = 2.85
for t,b,c in uvps:
    panel(s, ux, uy, uw, 2.28, accent=c, radius=0.05)
    tb(s, ux+0.45, uy+0.28, uw-0.75, 0.6, [{"runs":[(t, 17.5, True, WHITE, F_HEAD)]}])
    tb(s, ux+0.45, uy+1.0, uw-0.75, 1.2, [{"runs":[(b, 14, False, INK, F_BODY)],"ls":1.08}])
    uy += 2.46

# =====================================================================
# SLIDE 5 — TECHNICAL ARCHITECTURE
# =====================================================================
s = slide(); bg(s); header(s, "Technical Architecture", "Inputs → AI processing → confidence-anchored outputs.", 5)
colx = [0.7, 7.3, 14.0]; colw3 = [6.0, 6.0, 5.3]
heads = [("INPUTS", MUTE), ("AI PROCESSING", CYAN), ("OUTPUTS", MUTE)]
for i,(h,c) in enumerate(heads):
    tb(s, colx[i], 2.75, colw3[i], 0.4, [{"runs":[(h, 14, True, c, F_HEAD)],"align":PP_ALIGN.CENTER}])
# inputs
panel(s, colx[0], 3.25, colw3[0], 7.2, fill=PANEL, radius=0.04)
ins = [("Scenario","Natural-language query or map pin-drop"),
       ("Network","OpenStreetMap roads, jeepney & bike routes, heritage tags"),
       ("People","PSA 2020/2024 census, APIS poverty by barangay"),
       ("Environment","Sentinel-2 land cover, PAGASA / NOAH flood hazard"),
       ("Factors","WHO/EMEP emission factors, BIR zonal land values")]
iy = 3.5
for t,b in ins:
    tb(s, colx[0]+0.4, iy, colw3[0]-0.7, 1.3, [
        {"runs":[(t, 14.5, True, CYAN, F_HEAD)],"sa":2},
        {"runs":[(b, 13, False, INK, F_BODY)],"ls":1.04}])
    iy += 1.34
# processing (center, emphasized)
panel(s, colx[1], 3.25, colw3[1], 7.2, fill=C(0x0E2456), line=COBALT, lw=1.5, radius=0.04)
proc = [("Gemini 3.1 Orchestrator","NLP","parse query → simulation plan",COBALT),
        ("Unified Simulation Kernel","AGENTIC SIM","SUMO + persona pool + bias auditor → one trajectory dataset",CYAN),
        ("Five Impact Modules","ML","Behavioral · Social · Economic · Ecological · Societal, in parallel",CYAN),
        ("Synthesis + GraphRAG","GENAI + RAG","cited narrative; retrieval over the city knowledge graph",COBALT)]
py_ = 3.5
for t,tag,b,c in proc:
    panel(s, colx[1]+0.3, py_, colw3[1]-0.6, 1.55, fill=C(0x102C63), line=c, lw=1.0, radius=0.06)
    tb(s, colx[1]+0.55, py_+0.18, colw3[1]-1.1, 0.5, [{"runs":[(t, 15.5, True, WHITE, F_HEAD)]}])
    chip(s, colx[1]+colw3[1]-0.6-1.5, py_+0.2, 1.5, 0.38, tag, c, C(0x06163A), size=10.5)
    tb(s, colx[1]+0.55, py_+0.78, colw3[1]-1.1, 0.7, [{"runs":[(b, 12.5, False, INK, F_BODY)],"ls":1.04}])
    py_ += 1.67
# outputs
panel(s, colx[2], 3.25, colw3[2], 7.2, fill=PANEL, radius=0.04)
outs = [("Animated playback","Deck.gl TripsLayer — agents re-routing in real time",ENV),
        ("Five scored dimensions","each with a High / Medium / Low confidence band",COBALT),
        ("Glass-box Inspect","click any number → equation + datasets + confidence",CYAN),
        ("Exportable report","a structured recommendation for the planner",SOC)]
oy = 3.5
for t,b,c in outs:
    tb(s, colx[2]+0.4, oy, colw3[2]-0.7, 1.6, [
        {"runs":[(t, 14.5, True, c, F_HEAD)],"sa":2},
        {"runs":[(b, 13, False, INK, F_BODY)],"ls":1.04}])
    oy += 1.66
# arrows between columns
arrow(s, 6.80, 6.55, 0.40, 0.55, COBALT)
arrow(s, 13.52, 6.55, 0.40, 0.55, COBALT)
tb(s, 0.7, 10.6, 18.6, 0.5, [{"runs":[
    ("Engineered to a 90-second end-to-end budget: pre-warmed personas, delta simulation vs a nightly baseline, "
     "parallel modules, and a streaming progressive UI.", 13, False, MUTE, F_BODY, True)],"align":PP_ALIGN.CENTER}])

# =====================================================================
# SLIDE 6 — AI APPROACH & MODEL SELECTION
# =====================================================================
s = slide(); bg(s); header(s, "AI Approach & Model Selection", "Four AI techniques, each model chosen on purpose.", 6)
chips6 = ["GENERATIVE AI","AGENTIC SIMULATION","MACHINE LEARNING","RETRIEVAL-AUGMENTED GENERATION"]
cxx = 0.7
for c in chips6:
    w = 0.18*len(c)+0.5
    chip(s, cxx, 2.75, w, 0.5, c, PANEL2, CYAN, size=12, line=COBALT); cxx += w+0.3
models = [
 ("Gemini 3.1 Pro","Orchestration & synthesis","Parses plain language into a sim plan; writes the cited narrative. Current generation — 1.5 and 2.0 are shut down.",COBALT),
 ("Gemini 3.1 Flash-Lite","High-volume persona generation","200–500 commuter personas on the free tier — cheap and fast enough to pre-warm the pool.",CYAN),
 ("Eclipse SUMO (TraCI)","Physical multi-agent mobility","The open urban-mobility standard (DLR). Not OASIS / MiroFish — those model social media, not cities.",ENV),
 ("XGBoost","Corridor-volume baseline","Time-series forecaster for the nightly baseline that delta simulations run against.",ECO),
 ("bge-small-en","Knowledge-graph embeddings","Sentence-Transformers vectors in ChromaDB powering GraphRAG retrieval.",SOCI),
 ("Bias auditor (custom)","Mode-share anchoring","Anchors generated mode-share to ground truth and logs deviations to a public audit.",SOC),
]
cw6, g6, x06, y06 = 5.93, 0.4, 0.7, 3.55
for i,(name,role,desc,c) in enumerate(models):
    col = i % 3; rowi = i // 3
    x = x06 + col*(cw6+g6); y = y06 + rowi*2.55
    panel(s, x, y, cw6, 2.32, accent=c, radius=0.05)
    tb(s, x+0.42, y+0.26, cw6-0.7, 0.5, [{"runs":[(name, 17, True, WHITE, F_HEAD)]}])
    tb(s, x+0.42, y+0.78, cw6-0.7, 0.4, [{"runs":[(role.upper(), 11.5, True, c, F_HEAD)]}])
    tb(s, x+0.42, y+1.22, cw6-0.7, 1.0, [{"runs":[(desc, 13, False, INK, F_BODY)],"ls":1.05}])
panel(s, 0.7, 8.95, 18.6, 0.95, fill=C(0x081A45), line=COBALT, radius=0.07, accent=CYAN)
tb(s, 1.1, 9.08, 17.8, 0.7, [{"runs":[
   ("The guardrail:  ", 16.5, True, CYAN, F_HEAD),
   ("the LLM narrates and cites — it never originates a number. One kernel feeds five modules so the dimensions "
    "score one reality and cannot contradict each other.", 15.5, False, INK, F_BODY)]}], anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# SLIDE 7 — DATA STRATEGY
# =====================================================================
s = slide(); bg(s); header(s, "Data Strategy", "Open-data-first, contact-free, and licence-clean.", 7)
tiers = [
 ("TIER 1 · Day-1 open","OSM Geofabrik (ODbL) · PSA 2020/2024 census · Sentinel-2 (ESA) · PAGASA / Project NOAH flood · WHO/EMEP emission factors", ENV),
 ("TIER 2 · Request / FOI","LTFRB Region VI routes · PSA APIS poverty · Iloilo City CLUP 2021–2029 — filed Day 1; open substitutes used until they arrive", ECO),
 ("TIER 4 · Academic baseline","Calderon 2014 BRT corridor model · Macalalag 2021 bike study — anchors for validation", SOCI),
]
ty7 = 2.85
for t,b,c in tiers:
    panel(s, 0.7, ty7, 11.6, 1.62, accent=c, radius=0.05)
    tb(s, 1.15, ty7+0.22, 11.0, 0.4, [{"runs":[(t, 15.5, True, c, F_HEAD)]}])
    tb(s, 1.15, ty7+0.74, 11.0, 0.8, [{"runs":[(b, 13.5, False, INK, F_BODY)],"ls":1.05}])
    ty7 += 1.78
# licensing + cleaning
panel(s, 0.7, 8.25, 11.6, 2.4, fill=PANEL2, radius=0.05)
tb(s, 1.15, 8.45, 11.0, 0.4, [{"runs":[("LICENSING & CLEANING", 13.5, True, CYAN, F_HEAD)]}])
tb(s, 1.15, 8.95, 11.0, 1.7, [
   {"runs":[("Licences honoured: ", 14, True, WHITE, F_HEAD),("ODbL attribution, PSA & ESA open terms — no scraping behind blocks.", 14, False, INK, F_BODY)],"sa":6,"ls":1.05},
   {"runs":[("Reproducible build: ", 14, True, WHITE, F_HEAD),("idempotent fetch scripts; CCHAIN subset to Iloilo; BIR zonal values parsed; vintages stamped (newest first).", 14, False, INK, F_BODY)],"ls":1.05}])
# right: foundation stats
panel(s, 12.55, 2.85, 6.75, 7.8, fill=C(0x081A45), line=BORDER, radius=0.05)
tb(s, 12.95, 3.1, 6.0, 0.4, [{"runs":[("THE FOUNDATION IN HAND", 13.5, True, MUTE, F_HEAD)]}])
fstats = [("180","Iloilo barangays in the data foundation"),
          ("78 km²","pilot-city footprint — a full behavioral model"),
          ("5,680","BIR-priced land parcels parsed (RDO 74)"),
          ("2024","newest vintages preferred (POPCEN-CBMS)")]
fy = 3.7
for n,l in fstats:
    tb(s, 12.95, fy, 6.0, 0.8, [{"runs":[(n, 40, True, CYAN, F_NUM)]}], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 12.95, fy+0.82, 6.0, 0.6, [{"runs":[(l, 13, False, INK, F_BODY)],"ls":1.0}])
    fy += 1.72

# =====================================================================
# SLIDE 8 — AI ETHICS & RESPONSIBILITY
# =====================================================================
s = slide(); bg(s); header(s, "AI Ethics & Responsibility", "Honesty engineered in — not bolted on.", 8)
eth = [
 ("Bias, audited in the open","The bias auditor anchors generated persona mode-share to ground-truth mode-share and "
  "logs every deviation to a public audit. (Today it flags and surfaces; automated reweighting is on the roadmap — we say so.)", SOC),
 ("Confidence over false precision","Every dimension advertises a confidence floor and reports ranges, never a fake "
  "point estimate — the safeguard data-sparse ASEAN cities actually need.", COBALT),
 ("The informal economy, first-class","Jeepney, tricycle, and street-vendor displacement are modelled explicitly in "
  "Social and Economic — exactly what imported Western tools omit.", ECO),
 ("Privacy by design (RA 10173)","No personal data in the pipeline — open and aggregate sources only. The optional "
  "GPS-trace companion app is gated behind a Privacy Impact Assessment and counsel.", ENV),
]
cw8, g8, x08, y08 = 9.13, 0.34, 0.7, 2.95
for i,(t,b,c) in enumerate(eth):
    col=i%2; rowi=i//2
    x = x08+col*(cw8+g8); y = y08+rowi*2.85
    panel(s, x, y, cw8, 2.62, accent=c, radius=0.05)
    tb(s, x+0.45, y+0.3, cw8-0.8, 0.6, [{"runs":[(t, 18, True, WHITE, F_HEAD)]}])
    tb(s, x+0.45, y+1.05, cw8-0.8, 1.4, [{"runs":[(b, 14.5, False, INK, F_BODY)],"ls":1.08}])
panel(s, 0.7, 8.85, 18.6, 1.0, fill=C(0x081A45), line=COBALT, radius=0.07, accent=CYAN)
tb(s, 1.1, 9.0, 17.8, 0.7, [{"runs":[
   ("The glass box is the accountability mechanism:  ", 16.5, True, CYAN, F_HEAD),
   ("every number on screen is auditable to its equation, its open datasets, and a computed confidence.", 15.5, False, INK, F_BODY)]}], anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# SLIDE 9 — PROTOTYPE DEMONSTRATION
# =====================================================================
s = slide(); bg(s); header(s, "Prototype Demonstration", "It runs today — and it shows its work.", 9)
# two media placeholders
panel(s, 0.7, 2.9, 11.2, 6.0, fill=C(0x05122E), line=COBALT, lw=1.25, radius=0.03)
tb(s, 0.7, 5.2, 11.2, 1.2, [
   {"runs":[("▶  LIVE DEMO / GIF", 22, True, CYAN, F_HEAD)],"align":PP_ALIGN.CENTER,"sa":6},
   {"runs":[("Deck.gl TripsLayer — simulated commuters re-routing around the new project", 14, False, MUTE, F_BODY)],"align":PP_ALIGN.CENTER}])
tb(s, 1.1, 8.35, 10.4, 0.4, [{"runs":[("Drop  scenario-playback.png  /  recorded run here", 11.5, False, FAINT, F_BODY, True)],"align":PP_ALIGN.CENTER}])
panel(s, 12.25, 2.9, 7.05, 2.85, fill=C(0x05122E), line=CYAN, lw=1.25, radius=0.04)
tb(s, 12.25, 3.9, 7.05, 1.0, [
   {"runs":[("🔍  INSPECT DRAWER", 17, True, CYAN, F_HEAD)],"align":PP_ALIGN.CENTER,"sa":4},
   {"runs":[("equation → datasets → computed confidence", 12.5, False, MUTE, F_BODY)],"align":PP_ALIGN.CENTER}])
# the AI logic steps
panel(s, 12.25, 6.0, 7.05, 2.9, fill=PANEL, radius=0.05)
tb(s, 12.65, 6.2, 6.3, 0.4, [{"runs":[("THE AI “LOGIC”, ON SCREEN", 12.5, True, MUTE, F_HEAD)]}])
steps = [("1","Plain-language query → Gemini sim plan"),
         ("2","Kernel runs → map animates while modules compute"),
         ("3","Five dimension cards stream in with confidence"),
         ("4","Click a number → the glass box opens")]
sy = 6.7
for n,t in steps:
    chip(s, 12.65, sy, 0.42, 0.42, n, COBALT, WHITE, size=13)
    tb(s, 13.25, sy-0.02, 5.9, 0.5, [{"runs":[(t, 13.5, False, INK, F_BODY)]}], anchor=MSO_ANCHOR.MIDDLE)
    sy += 0.55
tb(s, 0.7, 9.15, 18.6, 0.7, [{"runs":[
   ("Live 90-second run in the pitch, with a recorded fallback ready. ", 14.5, True, WHITE, F_HEAD),
   ("The kernel, all five modules, the streaming API, and this frontend run end-to-end on 180 Iloilo barangays.",
    14.5, False, INK, F_BODY)]}])

# =====================================================================
# SLIDE 10 — TECHNICAL HURDLES
# =====================================================================
s = slide(); bg(s); header(s, "Technical Hurdles", "What broke — and how we made it honest.", 10)
hur = [
 ("Cross-dimension contradiction","Five separate simulators disagreed on the same project.",
  "One unified kernel → a single trajectory dataset every module scores. They physically cannot diverge.",CYAN),
 ("The model inventing numbers","An LLM will happily hallucinate a confident statistic.",
  "A glass-box contract + citation guard + an auditor that blocks the merge if any number ships unprovenanced.",COBALT),
 ("The 90-second budget","A cold run still clocks ~123 s, over target.",
  "Pre-warmed personas, delta-vs-baseline, parallel modules, and a trajectory cache — repeat runs return in <1 s.",ECO),
 ("Data gaps & provisional inputs","Mode-share isn’t calibrated; one flood fixture is provisional.",
  "We label them: Behavioral stays Medium, the fixture reads PROVISIONAL. Labelling provisional data is the feature.",SOC),
]
cw10,g10,x010,y010 = 9.13,0.34,0.7,2.95
for i,(t,prob,fix,c) in enumerate(hur):
    col=i%2; rowi=i//2
    x=x010+col*(cw10+g10); y=y010+rowi*2.95
    panel(s, x, y, cw10, 2.72, accent=c, radius=0.05)
    tb(s, x+0.45, y+0.26, cw10-0.8, 0.5, [{"runs":[(t, 17.5, True, WHITE, F_HEAD)]}])
    tb(s, x+0.45, y+0.86, cw10-0.8, 0.7, [{"runs":[("HURDLE  ", 11, True, SOC, F_HEAD),(prob, 13.5, False, MUTE, F_BODY)],"ls":1.04}])
    tb(s, x+0.45, y+1.6, cw10-0.8, 1.0, [{"runs":[("FIX  ", 11, True, GOOD, F_HEAD),(fix, 13.5, False, INK, F_BODY)],"ls":1.05}])
tb(s, 0.7, 9.05, 18.6, 0.7, [{"runs":[
   ("Every stage is now measured and visible — we optimise against real numbers, and we say “target,” "
    "“planned,” and “provisional” out loud. That discipline is the brand.", 14, False, MUTE, F_BODY, True)],"align":PP_ALIGN.CENTER}])

# =====================================================================
# SLIDE 11 — ACCURACY & EFFICIENCY METRICS
# =====================================================================
s = slide(); bg(s); header(s, "Accuracy & Efficiency", "Computed, not asserted — and honest about what isn’t ready.", 11)
# efficiency stat row
tb(s, 0.7, 2.7, 18.6, 0.4, [{"runs":[("EFFICIENCY — MEASURED", 13.5, True, CYAN, F_HEAD)]}])
estats = [("186","kernel tests pass (1 skipped)",GOOD),
          ("90 s","end-to-end budget (target)",COBALT),
          ("<1 s","warm / cached repeat run",CYAN),
          ("5","dimensions, one kernel",BEH),
          ("~123 s","cold run — over target, honest",WARN)]
ew, eg, ex = 3.48, 0.30, 0.7
for i,(n,l,c) in enumerate(estats):
    x = ex+i*(ew+eg)
    panel(s, x, 3.2, ew, 1.95, fill=PANEL, line=BORDER, radius=0.06)
    tb(s, x, 3.42, ew, 0.95, [{"runs":[(n, 44, True, c, F_NUM)]}], anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    tb(s, x+0.2, 4.38, ew-0.4, 0.7, [{"runs":[(l, 12.5, False, MUTE, F_BODY)],"align":PP_ALIGN.CENTER,"ls":1.0}])
# accuracy: validation harness
tb(s, 0.7, 5.45, 18.6, 0.4, [{"runs":[("ACCURACY — A VALIDATION HARNESS THAT GRADES ITSELF", 13.5, True, CYAN, F_HEAD)]}])
val = [("VAL-01 · Behavioral","Normalised RMSE vs the Calderon 2014 Iloilo BRT corridor.","Threshold ≤ 0.30  (FHWA corridor band)",COBALT),
       ("VAL-02 · Flood","Spatial IoU back-test vs the 2024 Iloilo flood extent.","Threshold ≥ 0.50  (Horritt & Bates 2002)",ENV)]
vx=0.7
for t,b,th,c in val:
    panel(s, vx, 5.95, 9.13, 1.95, accent=c, radius=0.05)
    tb(s, vx+0.45, 6.16, 8.3, 0.5, [{"runs":[(t, 16.5, True, WHITE, F_HEAD)]}])
    tb(s, vx+0.45, 6.7, 8.3, 0.6, [{"runs":[(b, 13.5, False, INK, F_BODY)],"ls":1.03}])
    tb(s, vx+0.45, 7.36, 8.3, 0.4, [{"runs":[(th, 13, True, c, F_NUM)]}])
    vx += 9.13+0.34
# honesty panel
panel(s, 0.7, 8.2, 18.6, 1.5, fill=C(0x2A1A06), line=C(0x7A5A1F), radius=0.06, accent=WARN)
tb(s, 1.15, 8.4, 17.9, 1.2, [{"runs":[
   ("The headline RMSE is deliberately withheld. ", 16, True, WARN, F_HEAD),
   ("Mode-share isn’t calibrated yet, so Behavioral stays Medium and the flood fixture is PROVISIONAL. Publishing "
    "a confident number from uncalibrated demand would violate the glass-box principle the product is built on. The "
    "harness is built, it runs, and it publishes the moment the calibration data lands.", 15, False, INK, F_BODY)],"ls":1.06}])

# =====================================================================
# SLIDE 12 — SCALABILITY ROADMAP
# =====================================================================
s = slide(); bg(s); header(s, "Scalability Roadmap", "One city to ASEAN — a config change, not a rebuild.", 12)
panel(s, 0.7, 2.85, 18.6, 1.5, fill=PANEL2, line=COBALT, radius=0.05)
tb(s, 1.15, 3.06, 17.9, 1.1, [
  {"runs":[("Zero hardware. ", 17, True, CYAN, F_HEAD),
           ("Scaling to a new city = swap the OpenStreetMap bounding box (API-level) and reweight the commuter "
            "personas to local modes (prompt-level). The city-agnostic CityConfig already lives in the code.", 16, False, INK, F_BODY)],"sa":4,"ls":1.05},
  {"runs":[("The cost of a new ASEAN city is API tokens — not procurement.", 15.5, True, WHITE, F_HEAD)]}])
steps = [("Iloilo","Pilot city","jeepney · tricycle · pedicab",ENV,"NOW"),
         ("PH cities","Cebu · Davao · Bacolod","same modes, new bbox",BEH,"NEXT"),
         ("Jakarta","Indonesia","ojek · angkot",ECO,"ASEAN"),
         ("Bangkok","Thailand","songthaew · tuk-tuk",SOCI,"ASEAN"),
         ("Ho Chi Minh / KL","Vietnam · Malaysia","xe-om · RapidKL",SOC,"ASEAN")]
sw5, sg, sx, syy = 3.48, 0.30, 0.7, 4.75
for i,(city,sub,modes,c,tag) in enumerate(steps):
    x = sx+i*(sw5+sg)
    panel(s, x, syy, sw5, 4.4, fill=PANEL, line=c, lw=1.25, radius=0.05)
    chip(s, x+0.3, syy+0.3, 1.5, 0.45, tag, c, C(0x06163A), size=12)
    tb(s, x+0.32, syy+1.05, sw5-0.6, 0.7, [{"runs":[(city, 21, True, WHITE, F_HEAD)],"ls":0.98}])
    tb(s, x+0.32, syy+1.95, sw5-0.6, 0.5, [{"runs":[(sub, 13.5, True, c, F_HEAD)]}])
    tb(s, x+0.32, syy+2.5, sw5-0.6, 1.6, [{"runs":[("Informal transit modes:", 12.5, False, MUTE, F_BODY)],"sa":3,"ls":1.0},
                                          {"runs":[(modes, 14.5, True, INK, F_HEAD)],"ls":1.05}])
    if i < len(steps)-1:
        arrow(s, x+sw5+0.03, syy+1.95, 0.24, 0.5, COBALT)
tb(s, 0.7, 9.4, 18.6, 0.5, [{"runs":[
   ("Naming the local informal-transit modes — ojek, angkot, songthaew, tuk-tuk, xe-om — is exactly what "
    "imported Western planning tools cannot do.", 13.5, False, MUTE, F_BODY, True)],"align":PP_ALIGN.CENTER}])

# =====================================================================
# SLIDE 13 — IMPACT ASSESSMENT
# =====================================================================
s = slide(); bg(s); header(s, "Impact Assessment", "Aligned to the UN SDGs and ASEAN resilience goals.", 13)
sdgs = [("11","Sustainable Cities",C(0xFD9D24)),
        ("13","Climate Action",C(0x3F7E44)),
        ("9","Industry & Innovation",C(0xFD6925)),
        ("10","Reduced Inequalities",C(0xDD1367)),
        ("3","Good Health",C(0x4C9F38))]
gw, gg, gx, gy = 3.48, 0.30, 0.7, 2.95
for i,(n,l,c) in enumerate(sdgs):
    x = gx+i*(gw+gg)
    panel(s, x, gy, gw, 2.05, fill=PANEL, line=BORDER, radius=0.06)
    sq = rect(s, x+0.35, gy+0.4, 1.25, 1.25, c, radius=0.08)
    tf=sq.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    rr=p.add_run(); rr.text="SDG\n"+n; rr.font.size=Pt(20); rr.font.bold=True; rr.font.name=F_HEAD; rr.font.color.rgb=WHITE
    tb(s, x+1.8, gy+0.5, gw-2.0, 1.1, [{"runs":[(l, 15, True, WHITE, F_HEAD)],"ls":1.0}], anchor=MSO_ANCHOR.MIDDLE)
# impact lines
impacts = [("De-risks capital before it’s committed","Flood-risk redistribution, climate exposure, and equity are scored pre-construction — not discovered on opening day.",COBALT),
           ("Equity is a first-class output","Equity-weighted access and informal-worker displacement are scored, not assumed (SDG 10).",SOC),
           ("A public-trust instrument","Confidence-anchored, auditable outputs are what data-sparse ASEAN cities need to trust a model.",ENV)]
iy13 = 5.4
for t,b,c in impacts:
    panel(s, 0.7, iy13, 12.4, 1.5, accent=c, radius=0.05)
    tb(s, 1.15, iy13+0.22, 11.7, 0.5, [{"runs":[(t, 16.5, True, WHITE, F_HEAD)]}])
    tb(s, 1.15, iy13+0.76, 11.7, 0.7, [{"runs":[(b, 13.5, False, INK, F_BODY)],"ls":1.04}])
    iy13 += 1.66
# regional anchor
panel(s, 13.3, 5.4, 6.0, 4.46, fill=C(0x081A45), line=COBALT, radius=0.05)
tb(s, 13.7, 5.65, 5.3, 0.4, [{"runs":[("REGIONAL ANCHOR", 13, True, MUTE, F_HEAD)]}])
tb(s, 13.7, 6.2, 5.3, 3.4, [
   {"runs":[("Iloilo City", 24, True, CYAN, F_HEAD)],"sa":4},
   {"runs":[("2026 ASEAN Clean Tourist City", 16, True, WHITE, F_HEAD)],"sa":10,"ls":1.05},
   {"runs":[("Aligns with the ASEAN Smart Cities Network and the AAIH 2026 sustainability theme. A representative "
             "suburban-metro testbed — not a special case.", 14, False, INK, F_BODY)],"ls":1.1}])

# =====================================================================
# SLIDE 14 — FUTURE ROADMAP
# =====================================================================
s = slide(); bg(s); header(s, "Future Roadmap", "After the hackathon: from prototype to pilot.", 14)
phases = [
 ("NOW — close the loop","Calibrate mode-share → publish the VAL-01 RMSE · replace the provisional flood fixture → publish VAL-02 IoU · drive the cold run under 90 s (libsumo / headless).",CYAN),
 ("NEXT — deepen the model","Bias auditor: flag → reweight · Hiligaynon gazetteer (colloquial place-name → GIS node) · CPDO feedback loop (PRD-F20) · RAG ingestion pipeline.",COBALT),
 ("PILOT — prove it in Iloilo","One real CPDO planner validates the demo → LGU + academic partnerships (Clean Air Asia, UP Visayas SURP) → a signed pilot scenario.",ENV),
 ("PRODUCT — sustain it","Public-good free tier for LGUs & academia; a paid private-developer / SaaS tier later. Deliberately TBD — adoption and credibility first.",ECO),
]
# timeline spine
rect(s, 1.2, 3.3, 0.07, 6.7, BORDER)
py14 = 3.05
for i,(t,b,c) in enumerate(phases):
    y = py14 + i*1.92
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y+0.18), Inches(0.55), Inches(0.55))
    dot.fill.solid(); dot.fill.fore_color.rgb=c; dot.line.color.rgb=WHITE; dot.line.width=Pt(1.5); dot.shadow.inherit=False
    panel(s, 1.9, y, 17.4, 1.66, accent=c, radius=0.04)
    tb(s, 2.35, y+0.26, 16.7, 0.5, [{"runs":[(t, 18, True, WHITE, F_HEAD)]}])
    tb(s, 2.35, y+0.84, 16.7, 0.7, [{"runs":[(b, 14.5, False, INK, F_BODY)],"ls":1.04}])

# =====================================================================
# SLIDE 15 — TEAM & CONTACT
# =====================================================================
s = slide(); bg(s, "bg_title.png")
rect(s, 0, 0, SW, 1.04, WHITE)
s.shapes.add_picture(A("logo_aaih.png"), Inches(0.7), Inches(0.17), height=Inches(0.72))
s.shapes.add_picture(A("logo_p2a_asean.png"), Inches(17.92), Inches(0.20), height=Inches(0.66))
tb(s, 0.7, 1.45, 18.6, 0.4, [{"runs":[("TEAM ATLAN  ·  POLYTECHNIC UNIVERSITY OF THE PHILIPPINES", 15, True, CYAN, F_HEAD)]}])
tb(s, 0.7, 1.85, 18.6, 0.9, [{"runs":[("The people behind MATRIX", 34, True, WHITE, F_HEAD)]}])
team = [
 ("Carlos Jerico Dela Torre","Team Lead","AI & Software Dev · Product & Business Architecture","CJ",COBALT),
 ("Yushin Bjorn Matsuda","Engineering","AI & Software Development · UI/UX Design","YM",CYAN),
 ("Maria Espina","Quality","QA · UI/UX Design","ME",ENV),
 ("Rica Mae Mago","Research","QA · Research & Marketing","RM",ECO),
 ("Russell Jay Fajardo","Research","QA · Research & Marketing","RF",SOCI),
]
tw5, tgap, tx5, tyy = 3.48, 0.30, 0.7, 3.15
for i,(name,role,desc,ini,c) in enumerate(team):
    x = tx5+i*(tw5+tgap)
    panel(s, x, tyy, tw5, 4.5, fill=PANEL, line=BORDER, radius=0.05)
    ov = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+tw5/2-0.75), Inches(tyy+0.4), Inches(1.5), Inches(1.5))
    ov.fill.solid(); ov.fill.fore_color.rgb=c; ov.line.color.rgb=WHITE; ov.line.width=Pt(1.5); ov.shadow.inherit=False
    tfo=ov.text_frame; tfo.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tfo.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    rr=p.add_run(); rr.text=ini; rr.font.size=Pt(30); rr.font.bold=True; rr.font.name=F_HEAD; rr.font.color.rgb=WHITE
    chip(s, x+tw5/2-1.0, tyy+2.12, 2.0, 0.42, role.upper(), c, C(0x06163A), size=11)
    tb(s, x+0.28, tyy+2.7, tw5-0.56, 0.9, [{"runs":[(name, 15.5, True, WHITE, F_HEAD)],"align":PP_ALIGN.CENTER,"ls":1.0}])
    tb(s, x+0.28, tyy+3.62, tw5-0.56, 0.8, [{"runs":[(desc, 12, False, MUTE, F_BODY)],"align":PP_ALIGN.CENTER,"ls":1.05}])
# contact bar
panel(s, 0.7, 8.1, 18.6, 1.5, fill=C(0x081A45), line=COBALT, radius=0.06)
tb(s, 1.0, 8.32, 18.0, 1.1, [
  {"runs":[("Contact:  ", 17, True, CYAN, F_HEAD),
           ("Carlos Jerico Dela Torre   ·   carlosjericodelatorre@gmail.com   ·   +63 949 636 9705", 17, False, WHITE, F_BODY)],"align":PP_ALIGN.CENTER,"sa":6},
  {"runs":[("MATRIX — decide before you build.", 17, True, CYAN, F_HEAD)],"align":PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
s.shapes.add_picture(A("logo_wordmark.png"), Inches(SW/2-1.6), Inches(11.7), height=Inches(0.42))

out = os.path.join(HERE, "Smart Cities_PUP_ATLAN_PitchDeck.pptx")
prs.save(out)
print("SAVED", out, "slides:", len(prs.slides._sldIdLst))
